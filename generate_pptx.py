from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Project ALO: Airtel Network Dashboard"
subtitle.text = "Network Intelligence & Management Vision\n"

# Slide 2: Executive Summary
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Executive Summary"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Comprehensive network visualization and management dashboard."
p = tf.add_paragraph()
p.text = "Real-time insights into network health and subscriber distribution."
p.level = 1
p = tf.add_paragraph()
p.text = "Monitors OLT cities, BRAS, and MSANs."
p.level = 1
p = tf.add_paragraph()
p.text = "Empowers operators to reduce downtime."
p.level = 1

# Slide 3: Core Features
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Core Features"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Geospatial Mapping: Interactive locations of hubs and OLTs."
p = tf.add_paragraph()
p.text = "Real-Time Alerting: Ranks OLT cities by complaint-to-subscriber ratios."
p = tf.add_paragraph()
p.text = "Analytics & Insights: Deep-dive charts and metrics on capacity."
p = tf.add_paragraph()
p.text = "Hierarchical Visualization: Drill-down from Circles to specific equipment."

# Slide 4: System Architecture
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "System Architecture & Data Flow"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Frontend: Next.js + React Context"
p = tf.add_paragraph()
p.text = "Map Layer: Leaflet for geospatial tracking"
p.level = 1
p = tf.add_paragraph()
p.text = "Charts Layer: Recharts for data analytics"
p.level = 1
p = tf.add_paragraph()
p.text = "Data Flow: Ingest -> Context -> Ratios -> UI Alert Generation"
p.level = 0

# Slide 5: Management Vision (AI Integration)
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Management Vision: AI Integration"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Goal: Proactive Network Intelligence Platform"
p = tf.add_paragraph()
p.text = "1. Predictive Maintenance: Anomaly detection on failure rates."
p.level = 1
p = tf.add_paragraph()
p.text = "2. Intelligent Routing: AI load forecasting and capacity planning."
p.level = 1
p = tf.add_paragraph()
p.text = "3. Automated Resolution (AIOps): Root Cause Analysis and Self-healing."
p.level = 1
p = tf.add_paragraph()
p.text = "Value: Reduce MTTR by up to 40%."
p.level = 1

prs.save('Project_ALO_Presentation.pptx')
print("Presentation generated successfully!")
